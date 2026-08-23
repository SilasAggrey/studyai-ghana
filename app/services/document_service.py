"""Document processing.

Extracts text from PDF (PyMuPDF), TXT/Markdown and DOCX (python-docx) so
uploaded lecture notes can be summarized, quizzed on, and interrogated.
"""
import logging

try:
    import pymupdf as fitz
except ImportError:  # older PyMuPDF releases
    import fitz  # type: ignore

from app.config import get_settings

logger = logging.getLogger(__name__)

CHUNK_SIZE = 1500
CHUNK_OVERLAP = 120

TEXT_FILE_TYPES = {"txt", "md", "markdown", "text"}


def extract_pdf_text(data: bytes, max_pages: int | None = None) -> tuple[str, int]:
    """Extract text from a PDF in memory. Returns (text, total_page_count)."""
    settings = get_settings()
    if max_pages is None:
        max_pages = settings.FREE_MAX_DOCUMENT_PAGES

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:
        raise ValueError(f"Could not open PDF: {exc}") from exc

    total_pages = doc.page_count
    pages = min(total_pages, max_pages)
    parts: list[str] = []
    for i in range(pages):
        try:
            parts.append(doc[i].get_text("text"))
        except Exception:
            logger.warning("failed to extract page %d", i)
    doc.close()
    return "\n".join(parts), total_pages


def extract_txt_text(data: bytes) -> tuple[str, int]:
    """Decode UTF-8 (with fallback) text files. Returns (text, 0 pages)."""
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding), 0
        except (UnicodeDecodeError, LookupError):
            continue
    raise ValueError("Could not read text file: unsupported encoding")


def extract_docx_text(data: bytes) -> tuple[str, int]:
    """Extract paragraphs + tables from a .docx in memory. Returns (text, 0)."""
    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(__import__("io").BytesIO(data))
    except Exception as exc:
        raise ValueError(f"Could not open DOCX: {exc}") from exc

    parts: list[str] = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            row_text = " | ".join(c for c in cells if c)
            if row_text:
                parts.append(row_text)
    return "\n".join(parts), 0


def extract_pptx_text(data: bytes) -> tuple[str, int]:
    """Extract slide titles + text + table contents from a .pptx. Returns (text, slide_count)."""
    try:
        from pptx import Presentation

        prs = Presentation(__import__("io").BytesIO(data))
    except Exception as exc:
        raise ValueError(f"Could not open PPTX: {exc}") from exc

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    row_text = " | ".join(c for c in cells if c)
                    if row_text:
                        slide_parts.append(row_text)
            if getattr(shape, "has_chart", False):
                try:
                    chart_texts = []
                    for series in shape.chart.plots[0].series:
                        if getattr(series, "name", None):
                            chart_texts.append(str(series.name))
                    for cat in getattr(shape.chart, "categories", None) or []:
                        if cat is not None:
                            chart_texts.append(str(cat))
                    if chart_texts:
                        slide_parts.append("Chart: " + ", ".join(chart_texts))
                except Exception:
                    pass
        if slide_parts:
            parts.append(f"[Slide {i}]")
            parts.extend(slide_parts)
    return "\n".join(parts), len(prs.slides)


def extract_text(data: bytes, file_type: str) -> tuple[str, int]:
    """Dispatch extraction by file extension. Returns (text, page_count)."""
    if file_type == "pdf":
        return extract_pdf_text(data)
    if file_type in TEXT_FILE_TYPES:
        return extract_txt_text(data)
    if file_type in ("docx",):
        return extract_docx_text(data)
    if file_type in ("pptx", "ppt"):
        return extract_pptx_text(data)
    raise ValueError(
        f"Unsupported file type: {file_type or 'unknown'}. Use PDF, TXT, DOCX, or PPTX."
    )


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks on paragraph boundaries where possible."""
    text = text.replace("\r\n", "\n").strip()
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        if end < len(text):
            # back off to a paragraph or sentence boundary
            boundary = max(
                text.rfind("\n\n", start, end),
                text.rfind(". ", start, end),
                text.rfind("\n", start, end),
            )
            if boundary > start + chunk_size // 2:
                end = boundary + 1
        chunks.append(text[start:end].strip())
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)
    return chunks


class DocumentService:
    def __init__(self, session):
        self.session = session

    async def store_document(
        self,
        user_id: int,
        *,
        filename: str,
        mime_type: str,
        telegram_file_id: str,
        size_bytes: int,
        data: bytes | None = None,
        title: str | None = None,
    ) -> object:
        from app.database.models import Document

        status = "uploaded"
        error = None
        page_count = 0
        text = None
        file_type = (
            filename.rsplit(".", 1)[-1].lower() if "." in filename else "unknown"
        )
        if data is not None:
            try:
                text, page_count = extract_text(data, file_type)
                status = "processed"
            except ValueError as exc:
                status = "failed"
                error = str(exc)[:250]

        doc = Document(
            user_id=user_id,
            filename=filename,
            file_type=file_type,
            mime_type=mime_type,
            telegram_file_id=telegram_file_id,
            size_bytes=size_bytes,
            page_count=page_count,
            status=status,
            title=title,
            extracted_text=text,
            error=error,
        )
        self.session.add(doc)
        await self.session.flush()
        return doc
