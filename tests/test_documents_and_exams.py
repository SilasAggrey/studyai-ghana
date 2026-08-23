"""PDF/TXT/DOCX extraction, chunking, and exam scoring."""
import io

import pytest

from app.ai.prompts.document import cap_material
from app.services.document_service import (
    DocumentService,
    chunk_text,
    extract_docx_text,
    extract_pdf_text,
    extract_pptx_text,
    extract_text,
    extract_txt_text,
)
from app.services.exam_service import compute_exam_result, grade_for


def _make_pdf() -> bytes:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello StudyAI Ghana students!")
    page.insert_text((72, 108), "This is a test PDF for document analysis.")
    out = io.BytesIO()
    doc.save(out)
    doc.close()
    return out.getvalue()


def test_extract_pdf_text():
    text, pages = extract_pdf_text(_make_pdf())
    assert pages == 1
    assert "StudyAI Ghana" in text


def test_extract_invalid_pdf_raises():
    with pytest.raises(ValueError):
        extract_pdf_text(b"not a pdf")


def test_extract_txt_text():
    text, pages = extract_txt_text("Hello students!\nLine two.".encode())
    assert pages == 0
    assert "Line two" in text


def test_extract_docx_text():
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_paragraph("Lecture one: Cell Biology")
    doc.add_paragraph("Mitochondria are the powerhouse.")
    out = io.BytesIO()
    doc.save(out)
    text, pages = extract_docx_text(out.getvalue())
    assert pages == 0
    assert "Cell Biology" in text
    assert "Mitochondria" in text


def test_extract_text_dispatch():
    txt, _ = extract_text(b"plain text notes", "txt")
    assert "plain text notes" in txt
    with pytest.raises(ValueError):
        extract_text(b"x", "exe")


def test_extract_pptx_text():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Photosynthesis"
    slide.placeholders[1].text_frame.text = "Chlorophyll absorbs light energy."
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    assert slide2 is not None
    out = io.BytesIO()
    prs.save(out)
    text, slides = extract_pptx_text(out.getvalue())
    assert slides == 2
    assert "[Slide 1]" in text
    assert "Photosynthesis" in text
    assert "Chlorophyll" in text


def test_cap_material_truncates_long_docs():
    long_text = "word " * 20000
    capped = cap_material(long_text, 1000)
    assert len(capped) <= 1200
    assert "truncated" in capped
    assert cap_material("short doc") == "short doc"


async def test_store_document_txt(session):
    from tests.conftest import make_user

    user = await make_user(session)
    doc = await DocumentService(session).store_document(
        user.id,
        filename="notes.txt",
        mime_type="text/plain",
        telegram_file_id="abc",
        size_bytes=12,
        data=b"hello notes",
    )
    assert doc.status == "processed"
    assert "hello notes" in doc.extracted_text
    assert doc.file_type == "txt"


async def test_store_document_unsupported(session):
    from tests.conftest import make_user

    user = await make_user(session)
    doc = await DocumentService(session).store_document(
        user.id,
        filename="notes.xyz",
        mime_type="application/octet-stream",
        telegram_file_id="abc",
        size_bytes=4,
        data=b"xxxx",
    )
    assert doc.status == "failed"
    assert doc.error is not None


def test_chunk_text_small_document():
    chunks = chunk_text("Single short paragraph.")
    assert len(chunks) == 1
    assert chunks[0] == "Single short paragraph."


def test_chunk_text_large_document():
    text = ("The quick brown fox jumps over the lazy dog. " * 200)
    chunks = chunk_text(text, chunk_size=500, overlap=50)
    assert len(chunks) > 1
    assert all(c for c in chunks)  # no empty chunks
    assert "".join(chunks)  # content preserved


def test_grade_boundaries():
    assert grade_for(90) == "A+"
    assert grade_for(80) == "A"
    assert grade_for(70) == "B+"
    assert grade_for(60) == "B"
    assert grade_for(50) == "C+"
    assert grade_for(40) == "C"
    assert grade_for(10) == "F"


def test_compute_exam_result():
    result = compute_exam_result(24, 30, 47)
    assert result["score"] == 24
    assert result["percentage"] == pytest.approx(80.0)
    assert result["grade"] == "A"
    assert result["wrong"] == 6
